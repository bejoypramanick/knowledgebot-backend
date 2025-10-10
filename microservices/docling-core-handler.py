#!/usr/bin/env python3
"""
Docling Core - Lightweight Document Processing
Handles simple documents without heavy dependencies
Optimized for text extraction and basic processing
"""

import json
import os
import base64
from typing import Dict, Any, List
import logging
import io
from pathlib import Path

logger = logging.getLogger(__name__)
logger.setLevel(logging.INFO)

# Initialize lightweight libraries
try:
    import pypdfium2 as pdfium
    import pdfplumber
    from docx import Document
    from pptx import Presentation
    import pandas as pd
    logger.info("Docling Core libraries initialized successfully")
except Exception as e:
    logger.error(f"Failed to initialize Docling Core libraries: {e}")
    pdfium = None
    pdfplumber = None
    Document = None
    Presentation = None
    pd = None

def extract_text_from_pdf(document_bytes: bytes) -> Dict[str, Any]:
    """Extract text from PDF using lightweight methods"""
    results = {
        "success": False,
        "text": "",
        "pages": [],
        "metadata": {},
        "method_used": None
    }
    
    try:
        # Method 1: PyPDFium2 (fastest)
        if pdfium:
            pdf = pdfium.PdfDocument(document_bytes)
            text_content = ""
            pages = []
            
            for i in range(len(pdf)):
                page = pdf[i]
                textpage = page.get_textpage()
                text = textpage.get_text_range()
                text_content += text + "\n"
                pages.append({
                    "page_number": i + 1,
                    "text": text,
                    "char_count": len(text)
                })
            
            results.update({
                "success": True,
                "text": text_content,
                "pages": pages,
                "method_used": "pypdfium2",
                "metadata": {
                    "total_pages": len(pdf),
                    "total_text_length": len(text_content)
                }
            })
            
        # Method 2: PDFPlumber (fallback)
        elif pdfplumber:
            with pdfplumber.open(io.BytesIO(document_bytes)) as pdf:
                text_content = ""
                pages = []
                
                for i, page in enumerate(pdf.pages):
                    text = page.extract_text() or ""
                    text_content += text + "\n"
                    pages.append({
                        "page_number": i + 1,
                        "text": text,
                        "char_count": len(text)
                    })
                
                results.update({
                    "success": True,
                    "text": text_content,
                    "pages": pages,
                    "method_used": "pdfplumber",
                    "metadata": {
                        "total_pages": len(pdf.pages),
                        "total_text_length": len(text_content)
                    }
                })
        
        else:
            results["error"] = "No PDF processing libraries available"
            
    except Exception as e:
        logger.error(f"PDF text extraction failed: {e}")
        results["error"] = str(e)
    
    return results

def extract_text_from_docx(document_bytes: bytes) -> Dict[str, Any]:
    """Extract text from DOCX files"""
    results = {
        "success": False,
        "text": "",
        "paragraphs": [],
        "tables": [],
        "metadata": {}
    }
    
    try:
        if Document:
            doc = Document(io.BytesIO(document_bytes))
            
            # Extract paragraphs
            paragraphs = []
            full_text = ""
            
            for i, paragraph in enumerate(doc.paragraphs):
                text = paragraph.text.strip()
                if text:
                    paragraphs.append({
                        "index": i,
                        "text": text,
                        "style": paragraph.style.name if paragraph.style else "Normal"
                    })
                    full_text += text + "\n"
            
            # Extract tables
            tables = []
            for i, table in enumerate(doc.tables):
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                
                tables.append({
                    "index": i,
                    "data": table_data,
                    "rows": len(table_data),
                    "cols": len(table_data[0]) if table_data else 0
                })
            
            results.update({
                "success": True,
                "text": full_text,
                "paragraphs": paragraphs,
                "tables": tables,
                "metadata": {
                    "total_paragraphs": len(paragraphs),
                    "total_tables": len(tables),
                    "total_text_length": len(full_text)
                }
            })
        else:
            results["error"] = "DOCX processing library not available"
            
    except Exception as e:
        logger.error(f"DOCX text extraction failed: {e}")
        results["error"] = str(e)
    
    return results

def extract_text_from_pptx(document_bytes: bytes) -> Dict[str, Any]:
    """Extract text from PPTX files"""
    results = {
        "success": False,
        "text": "",
        "slides": [],
        "metadata": {}
    }
    
    try:
        if Presentation:
            prs = Presentation(io.BytesIO(document_bytes))
            
            slides = []
            full_text = ""
            
            for i, slide in enumerate(prs.slides):
                slide_text = ""
                slide_elements = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        slide_text += text + "\n"
                        slide_elements.append({
                            "type": shape.shape_type,
                            "text": text
                        })
                
                slides.append({
                    "slide_number": i + 1,
                    "text": slide_text,
                    "elements": slide_elements,
                    "char_count": len(slide_text)
                })
                full_text += slide_text + "\n"
            
            results.update({
                "success": True,
                "text": full_text,
                "slides": slides,
                "metadata": {
                    "total_slides": len(slides),
                    "total_text_length": len(full_text)
                }
            })
        else:
            results["error"] = "PPTX processing library not available"
            
    except Exception as e:
        logger.error(f"PPTX text extraction failed: {e}")
        results["error"] = str(e)
    
    return results

def extract_text_from_txt(document_bytes: bytes) -> Dict[str, Any]:
    """Extract text from plain text files"""
    try:
        # Try different encodings
        encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        text = ""
        
        for encoding in encodings:
            try:
                text = document_bytes.decode(encoding)
                break
            except UnicodeDecodeError:
                continue
        
        if not text:
            text = document_bytes.decode('utf-8', errors='ignore')
        
        return {
            "success": True,
            "text": text,
            "metadata": {
                "total_text_length": len(text),
                "encoding": "detected"
            }
        }
        
    except Exception as e:
        logger.error(f"Text extraction failed: {e}")
        return {
            "success": False,
            "error": str(e)
        }

def process_document_core(document_bytes: bytes, filename: str, document_id: str = None) -> Dict[str, Any]:
    """
    Process document using lightweight methods
    Optimized for simple text extraction without heavy dependencies
    """
    try:
        file_extension = Path(filename).suffix.lower()
        
        # Route to appropriate extractor
        if file_extension == '.pdf':
            result = extract_text_from_pdf(document_bytes)
        elif file_extension in ['.docx', '.doc']:
            result = extract_text_from_docx(document_bytes)
        elif file_extension == '.pptx':
            result = extract_text_from_pptx(document_bytes)
        elif file_extension in ['.txt', '.md', '.csv']:
            result = extract_text_from_txt(document_bytes)
        else:
            # Try text extraction as fallback
            result = extract_text_from_txt(document_bytes)
        
        # Add common metadata
        result.update({
            "document_id": document_id,
            "filename": filename,
            "file_extension": file_extension,
            "processing_service": "docling-core",
            "processing_time": 0  # Will be set by caller
        })
        
        return result
        
    except Exception as e:
        logger.error(f"Core document processing failed: {e}")
        return {
            "success": False,
            "error": f"Core processing failed: {str(e)}",
            "document_id": document_id,
            "filename": filename,
            "processing_service": "docling-core"
        }

def lambda_handler(event: Dict[str, Any], context: Any) -> Dict[str, Any]:
    """
    Lambda handler for Docling Core processing
    """
    try:
        # Extract document data
        document_bytes = base64.b64decode(event.get('document_bytes', ''))
        filename = event.get('filename', 'unknown')
        document_id = event.get('document_id')
        
        if not document_bytes:
            return {
                "statusCode": 400,
                "body": json.dumps({
                    "success": False,
                    "error": "No document bytes provided"
                })
            }
        
        # Process document
        import time
        start_time = time.time()
        result = process_document_core(document_bytes, filename, document_id)
        result["processing_time"] = time.time() - start_time
        
        return {
            "statusCode": 200,
            "body": json.dumps(result)
        }
        
    except Exception as e:
        logger.error(f"Lambda handler error: {e}")
        return {
            "statusCode": 500,
            "body": json.dumps({
                "success": False,
                "error": f"Core processing failed: {str(e)}"
            })
        }
